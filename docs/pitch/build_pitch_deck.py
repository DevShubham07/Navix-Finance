#!/usr/bin/env python3
"""Builds the SoftSolutionsAI x DhanBoost sales deck (.pptx).

    python3 docs/pitch/build_pitch_deck.py

Drop screenshots into docs/pitch/shots/<name>.png (names listed in SHOTLIST.md) and re-run —
the same slides pick them up in place of the dashed placeholder frames.

Re-skin by editing PALETTE / FONTS / LOGO below. Slide copy is the SLIDES list; nothing else
needs touching to change what the deck says.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
REPO = HERE.parent.parent
SHOTS = HERE / "shots"
OUT = HERE / "SoftSolutionsAI_DhanBoost_Pitch.pptx"
LOGO = REPO / "frontend" / "public" / "dhanboost-mark.png"

PALETTE = {
    "navy": RGBColor(0x0C, 0x25, 0x40),
    "gold": RGBColor(0xE9, 0xB5, 0x3A),
    "cream": RGBColor(0xFD, 0xFB, 0xF6),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "ink": RGBColor(0x10, 0x15, 0x1B),
    "muted": RGBColor(0x5B, 0x6B, 0x7C),
    "line": RGBColor(0xD8, 0xD2, 0xC4),
}
FONTS = {"body": "Verdana"}  # safe on Windows + macOS PowerPoint without embedding

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.72)  # page margin

# every textbox we create, for the overflow estimate in the self-check
OVERFLOW: list = []


# ---------------------------------------------------------------- primitives


def est_lines(body, width_in, pt):
    """Rough wrapped-line count. Nothing here renders the deck, so this is how we catch text
    that would spill off a slide. 0.52 = average glyph width as a fraction of point size."""
    if width_in <= 0:
        return 1
    per_line = max(8, int(width_in * 72 / (pt * (0.48 if pt >= 24 else 0.52))))
    return sum(max(1, -(-len(seg) // per_line)) for seg in body.split("\n"))


def blank(prs, bg="cream"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid()
    r.fill.fore_color.rgb = PALETTE[bg]
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(slide, x, y, w, h, fill=None, line=None, dash=False, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    if radius:
        shape.adjustments[0] = radius
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = PALETTE[fill]
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = PALETTE[line]
        shape.line.width = Pt(1.25)
        if dash:
            shape.line.dash_style = 4  # msoLineDash
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines: list of (string, size_pt, color_key, bold, space_after_pt) or plain str."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    expanded = []
    for spec in lines:
        if isinstance(spec, str):
            spec = (spec, 14, "ink", False, 8)
        body, size, color, bold, after = spec
        parts = body.split("\n")
        for j, part in enumerate(parts):  # a literal newline becomes its own paragraph
            expanded.append((part, size, color, bold, after if j == len(parts) - 1 else 0))
    for i, (body, size, color, bold, after) in enumerate(expanded):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(after)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = body
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = PALETTE[color]
        run.font.name = FONTS["body"]
    OVERFLOW.append((tb, expanded, w, h))
    return tb


def kicker(slide, label, color="gold"):
    text(slide, M, Inches(0.5), Inches(8), Inches(0.3), [(label.upper(), 11, color, True, 0)])


def heading(slide, title, sub=None, dark=False):
    """Title, optional subtitle, gold rule. Grows downward when either wraps, so a long
    title can never sit on top of its own subtitle."""
    kc = "white" if dark else "navy"
    tw = W - 2 * M
    th = Inches(est_lines(title, tw / Inches(1), 30) * 0.5)
    text(slide, M, Inches(0.82), tw, th, [(title, 30, kc, True, 0)])
    y = Inches(0.82) + th + Inches(0.1)
    if sub:
        sh = Inches(est_lines(sub, tw / Inches(1), 14) * 0.26)
        text(slide, M, y, tw, sh, [(sub, 14, "gold" if dark else "muted", False, 0)])
        y = y + sh + Inches(0.16)
    box(slide, M, y, Inches(1.4), Pt(3), fill="gold")
    return y + Inches(0.3)


def bullets(slide, x, y, w, items, size=12.5, gap=10):
    lines = []
    for it in items:
        if it.startswith("## "):
            lines.append((it[3:], size + 1, "navy", True, 5))
        elif it.startswith("> "):
            lines.append((it[2:], size - 2, "muted", False, gap))
        else:
            lines.append(("•  " + it, size, "ink", False, gap))
    return text(slide, x, y, w, H - y - M, lines)


def picture_or_frame(slide, name, caption, x, y, w, h):
    """Embed shots/<name>.png fitted into the box, else a dashed placeholder."""
    box(slide, x, y, w, h, fill="white", line="line")
    src = SHOTS / f"{name}.png"
    if src.exists():
        pic = slide.shapes.add_picture(str(src), x, y, width=w)
        if pic.height > h:
            pic.height, pic.width = h, Emu(int(pic.width * h / pic.height))
        pic.left = Emu(int(x + (w - pic.width) / 2))
        pic.top = Emu(int(y + (h - pic.height) / 2))
        return True
    inner = Inches(0.22)
    box(slide, x + inner, y + inner, w - 2 * inner, h - 2 * inner, line="gold", dash=True)
    text(
        slide,
        x + Inches(0.5),
        y + h / 2 - Inches(0.55),
        w - Inches(1.0),
        Inches(1.1),
        [("SCREENSHOT", 11, "gold", True, 6), (caption, 13, "muted", False, 0)],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return False


def logo(slide, size=Inches(0.55)):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), W - M - size, H - M - size, height=size)


# ---------------------------------------------------------------- slide kinds


def s_title(prs, d):
    s = blank(prs, "navy")
    box(s, 0, H - Inches(0.28), W, Inches(0.28), fill="gold")
    text(s, M, Inches(1.5), W - 2 * M, Inches(0.4), [(d["kicker"], 12, "gold", True, 0)])
    text(s, M, Inches(2.1), Inches(10.6), Inches(2.4), [(d["title"], 44, "white", True, 10)])
    text(s, M, Inches(4.35), Inches(9.6), Inches(1.6), [(d["sub"], 17, "cream", False, 0)])
    text(s, M, Inches(6.15), Inches(9.6), Inches(0.5), [(d["foot"], 12, "gold", False, 0)])
    logo(s, Inches(0.9))
    return s


def s_section(prs, d):
    s = blank(prs, "navy")
    text(s, M, Inches(2.4), Inches(2.2), Inches(1.6), [(d["num"], 68, "gold", True, 0)])
    text(s, M, Inches(3.6), Inches(10.4), Inches(1.2), [(d["title"], 34, "white", True, 8)])
    text(s, M, Inches(4.6), Inches(9.8), Inches(1.0), [(d["sub"], 15, "cream", False, 0)])
    return s


def s_bullets(prs, d):
    s = blank(prs)
    kicker(s, d.get("kicker", "SoftSolutionsAI"))
    y = heading(s, d["title"], d.get("sub"))
    cols = d.get("cols")
    if cols:
        gw = (W - 2 * M - Inches(0.6)) / 2
        bullets(s, M, y, gw, cols[0])
        bullets(s, M + gw + Inches(0.6), y, gw, cols[1])
    else:
        bullets(s, M, y, W - 2 * M, d["body"])
    logo(s)
    return s


def s_shot(prs, d):
    s = blank(prs)
    kicker(s, d.get("kicker", "product walkthrough"))
    y = heading(s, d["title"], d.get("sub"))
    fw, fh = Inches(7.11), Inches(4.0)  # 16:9, so a 1920x1080 grab drops in unscaled
    ok = picture_or_frame(s, d["shot"], d["frame"], M, y, fw, fh)
    gx = M + fw + Inches(0.4)
    bullets(s, gx, y - Inches(0.06), W - M - gx, d["body"], size=12, gap=9)
    text(s, M, y + fh + Inches(0.14), fw, Inches(0.3), [(d["route"], 10.5, "muted", False, 0)])
    d["_shot_ok"] = ok
    return s


def s_table(prs, d):
    s = blank(prs)
    kicker(s, d.get("kicker", "SoftSolutionsAI"))
    y = heading(s, d["title"], d.get("sub"))
    rows, cols = len(d["rows"]) + 1, len(d["head"])
    fs = d.get("fsize", 11)
    rh = Inches(d.get("rowh", 0.34))
    gt = s.shapes.add_table(rows, cols, M, y, W - 2 * M, rh * rows).table
    for i, wfrac in enumerate(d.get("widths", [1 / cols] * cols)):
        gt.columns[i].width = Emu(int((W - 2 * M) * wfrac))
    widths = [gt.columns[i].width / Inches(1) for i in range(cols)]
    total_h = Emu(0)
    for r, row in enumerate([d["head"]] + d["rows"]):
        # PowerPoint grows rows to fit text; reserve that height ourselves so the footnote
        # below the table lands under it instead of on top of it
        lines = max(est_lines(str(v), widths[c] - 0.18, fs) for c, v in enumerate(row))
        gt.rows[r].height = max(rh, Inches(lines * fs * 1.3 / 72 + 0.1))
        total_h = Emu(total_h + gt.rows[r].height)
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.text = ""
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PALETTE["navy" if r == 0 else "white"]
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(fs)
            run.font.bold = r == 0
            run.font.name = FONTS["body"]
            run.font.color.rgb = PALETTE["white" if r == 0 else "ink"]
    if d.get("foot"):
        fy = y + total_h + Inches(0.18)
        text(s, M, fy, W - 2 * M, H - fy - Inches(0.2), [(d["foot"], 11, "muted", False, 0)])
    return s


def s_chain(prs, d):
    """A wrapped row of boxes with arrows — journey maps / state machines."""
    s = blank(prs)
    kicker(s, d.get("kicker", "product walkthrough"))
    y0 = heading(s, d["title"], d.get("sub"))
    per = d.get("per_row", 5)
    steps = d["steps"]
    bw = (W - 2 * M - Inches(0.34) * (per - 1)) / per
    bh = Inches(0.98)
    for i, (label, note) in enumerate(steps):
        r, c = divmod(i, per)
        x = M + c * (bw + Inches(0.34))
        y = y0 + r * (bh + Inches(0.62))
        accent = label.startswith("*")
        box(s, x, y, bw, bh, fill="gold" if accent else "white", line="line", radius=0.14)
        text(
            s,
            x + Inches(0.12),
            y + Inches(0.11),
            bw - Inches(0.24),
            bh - Inches(0.2),
            [(label.lstrip("*"), 12, "navy", True, 3), (note, 9.5, "muted" if not accent else "navy", False, 0)],
        )
        if c < per - 1 and i < len(steps) - 1:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw + Inches(0.05), y + bh / 2 - Inches(0.08), Inches(0.24), Inches(0.16))
            a.fill.solid()
            a.fill.fore_color.rgb = PALETTE["gold"]
            a.line.fill.background()
            a.shadow.inherit = False
    if d.get("foot"):
        fy = y0 + (-(-len(steps) // per)) * (bh + Inches(0.62)) + Inches(0.1)
        text(s, M, fy, W - 2 * M, H - fy - Inches(0.2), [(d["foot"], 11.5, "muted", False, 0)])
    return s


def s_stats(prs, d):
    s = blank(prs)
    kicker(s, d.get("kicker", "SoftSolutionsAI"))
    y = heading(s, d["title"], d.get("sub"))
    per = 4
    bw = (W - 2 * M - Inches(0.3) * (per - 1)) / per
    bh = Inches(1.62)
    for i, (big, small) in enumerate(d["stats"]):
        r, c = divmod(i, per)
        x, yy = M + c * (bw + Inches(0.3)), y + r * (bh + Inches(0.22))
        box(s, x, yy, bw, bh, fill="white", line="line", radius=0.1)
        text(s, x + Inches(0.18), yy + Inches(0.16), bw - Inches(0.36), bh - Inches(0.3),
             [(big, 22, "navy", True, 4), (small, 10.5, "muted", False, 0)])
    if d.get("foot"):
        fy = y + 2 * bh + Inches(0.46)
        text(s, M, fy, W - 2 * M, H - fy - Inches(0.2), [(d["foot"], 11, "muted", False, 0)])
    logo(s)
    return s


def s_close(prs, d):
    s = blank(prs, "navy")
    box(s, 0, 0, Inches(0.28), H, fill="gold")
    x = M + Inches(0.3)
    text(s, x, Inches(1.45), Inches(10.4), Inches(1.9), [(d["title"], 33, "white", True, 0)])
    text(s, x, Inches(3.35), Inches(10.2), Inches(2.5), [(b, 13.5, "cream", False, 11) for b in d["body"]])
    text(s, x, Inches(6.3), Inches(10.2), Inches(0.5), [(d["foot"], 13, "gold", True, 0)])
    logo(s, Inches(0.8))
    return s


BUILDERS = {
    "title": s_title, "section": s_section, "bullets": s_bullets, "shot": s_shot,
    "table": s_table, "chain": s_chain, "stats": s_stats, "close": s_close,
}

# ---------------------------------------------------------------- deck content

SLIDES = [
    {
        "kind": "title",
        "kicker": "SoftSolutionsAI  ·  softsolutionsai.com",
        "title": "We ship lending platforms,\nnot prototypes.",
        "sub": "Case study: DhanBoost — a salary-linked digital lending platform, built end to end and "
               "running in production on AWS.",
        "foot": "Product walkthrough  ·  borrower journey + 9-role back office  ·  ~35 minutes",
        "notes": "Open with the promise, not the tech: 'Everything you're about to see is running "
                 "software — a real state machine, real KYC providers, a real audit trail. We built it "
                 "for one NBFC and we can build it for you.' Say the runtime up front so nobody "
                 "wonders how long this will take. Do NOT open the app yet — the deck sets up what "
                 "they should be looking for before they see a screen.",
    },
    {
        "kind": "bullets",
        "kicker": "the problem",
        "title": "Standing up a lending business is an integration problem",
        "sub": "Most vendors sell you one box on the diagram. You own everything between the boxes.",
        "cols": [[
            "## What you can buy off the shelf",
            "A LOS for origination. An LMS for the ledger. A collections tool. Three KYC vendors.",
            "Each one owns its own database, its own idea of a 'customer', and its own definition of 'approved'.",
            "## What you still have to build",
            "The state machine that decides what is allowed to happen next, and by whom.",
            "Separation of duties that an auditor will accept.",
            "One audit trail that spans origination, disbursal, repayment and collections.",
        ], [
            "## What it costs you",
            "12–18 months before the first rupee is disbursed.",
            "A permanent reconciliation tax between systems that each think they are the source of truth.",
            "A compliance story assembled after the fact from vendor exports.",
            "## The alternative",
            "One aggregate, one status field, one append-only event log — origination through "
            "write-off. That is what DhanBoost is, and it is what the rest of this deck walks through.",
        ]],
        "notes": "Land the pain before the product. If the room is a startup NBFC, the line that "
                 "works is '12-18 months before the first disbursal.' If it's an existing lender, the "
                 "line is 'the reconciliation tax' — they already live it. Ask which one they are and "
                 "lean on that half. Don't name competitor products.",
    },
    {
        "kind": "stats",
        "kicker": "what we delivered",
        "title": "DhanBoost, by the numbers",
        "sub": "Every figure below is countable in the repository — nothing here is a marketing estimate.",
        "stats": [
            ("1 aggregate", "One loan_application row walks the whole lifecycle. No stage-skipping."),
            ("9 roles", "KYC, credit exec/head, disbursement, accounting, collections x2, admin, developer."),
            ("8 checks", "PAN · email · address · Aadhaar · bureau · salary · penny-drop · selfie — all mandatory."),
            ("40 migrations", "Versioned Flyway schema, applied on every boot. No manual DDL."),
            ("51 test classes", "380+ backend test methods, plus a Testcontainers integration suite."),
            ("7 e2e specs", "Playwright: auth, RBAC, onboarding, repay, staff console, admin."),
            ("Live on AWS", "Vercel » ALB » ECS Fargate » RDS/S3/SSM. Documented redeploy runbook."),
            ("49 seeded apps", "A one-command demo environment at every lifecycle stage, shipped with it."),
        ],
        "foot": "Built by SoftSolutionsAI. The same spine is the starting point for your platform — "
                "the walkthrough that follows is what you would be buying, screen by screen.",
        "notes": "This is the credibility slide. Read three numbers, not eight — '8 mandatory checks, "
                 "51 test classes, live on AWS' — and move. If a technical buyer stops you here, the "
                 "detail is on slides 21-24; say so and keep the walkthrough moving. The 49-seeded-apps "
                 "number matters more than it looks: it's why we can demo every edge state on demand.",
    },
    {
        "kind": "table",
        "kicker": "the product",
        "title": "One product rule set, priced the same for everyone",
        "sub": "A salaried borrower draws a short advance, pays the fee upfront, and repays once on salary day.",
        "head": ["Rule", "Value", "Why it is built this way"],
        "widths": [0.24, 0.28, 0.48],
        "rows": [
            ["Eligible limit", "25% of monthly salary, floored to ₹100", "Affordability is derived from verified salary, not self-declared."],
            ["Minimum loan", "₹1,000", "Below this the unit economics don't clear the verification cost."],
            ["Processing fee", "10% of principal, upfront", "Deducted from disbursal — the borrower never owes it later."],
            ["GST", "18% on the fee, upfront", "Collected and shown separately on every quote."],
            ["Interest", "1% per day, over actual tenure", "Prepay early and you pay interest only to the day you pay."],
            ["Due date", "Next salary credit, within 40 days", "Repayment lands when money does. This is the core risk control."],
            ["Late penalty", "2% per day, capped at 30 days", "The cap is hard-coded — the balance cannot run away."],
            ["Repayment", "One single instalment", "No EMI schedule to service, dispute or reconcile."],
        ],
        "foot": "Worked example — ₹10,000 principal: fee ₹1,000 + GST ₹180 » borrower receives ₹8,820. "
                "Disbursed 3 June, salary day 30 » due 30 June (27 days) » repays ₹12,700. All money is "
                "integer paise, HALF_UP. Risk bands A/B/C/D change the limit and the checks, never the price.",
        "notes": "The last line is the one to say out loud: risk bands change the limit, never the "
                 "price. It pre-empts the 'is this risk-based pricing?' question and it's a real "
                 "product decision the client made. Also flag integer paise — anyone who has debugged "
                 "a floating-point ledger will react to that. Walk the ₹10,000 example on your fingers; "
                 "it makes the whole economics concrete in fifteen seconds.",
    },
    {
        "kind": "chain",
        "kicker": "borrower journey",
        "title": "Nine verified steps, then money",
        "sub": "Every step is a real provider call, stored with its provider reference. None of it is a form that says 'verified'.",
        "per_row": 5,
        "steps": [
            ("PAN", "Compliance search"), ("Email", "Deliverability"), ("Address", "Verification API"),
            ("DigiLocker", "Aadhaar, with consent"), ("Bureau", "Experian / CRIF"),
            ("Salary", "Income check"), ("Penny-drop", "Bank account owner"), ("Selfie", "Liveness + face match"),
            ("Agreement", "e-Signed, stored"), ("*Limit issued", "Apply, disburse, repay"),
        ],
        "foot": "The gate: submit-kyc is refused with KYC_INCOMPLETE until all eight checks clear. A borrower "
                "cannot talk their way into the credit queue, and staff cannot pull them in early.",
        "notes": "Point at the gate line, not the boxes. The differentiator versus a form-driven "
                 "onboarding is that completeness is enforced server-side — the API refuses. Mention "
                 "that a KYC approver can still manually override a single borderline check, and that "
                 "the override is itself audited (slide 14) — so the gate is strict without being "
                 "operationally stupid.",
    },
    {
        "kind": "shot",
        "kicker": "borrower journey  ·  1 of 5",
        "title": "Sign-up: mobile OTP, PAN, email, address",
        "sub": "Four steps, four real API calls, resumable if the borrower drops off.",
        "shot": "borrower-onboarding",
        "frame": "Sign-up wizard — mobile OTP » PAN » email » address",
        "route": "/signup/mobile-otp  ·  /signup/pan  ·  /signup/email  ·  /signup/address",
        "body": [
            "OTP over a DLT-registered SMS template — the session is established at step one, so a "
            "drop-off resumes where it left off.",
            "PAN runs a compliance search, not a checksum. The name that comes back is what the rest "
            "of the file is checked against.",
            "Email is deliverability-verified, which is what makes the email channel safe to send to later.",
            "Address is verified through a provider API — no manual data entry to reconcile.",
            "> Each result is stored as its own verification row with the provider and its reference, "
            "so any check can be re-run or explained a year later.",
        ],
        "notes": "Keep this one short — it's table stakes and the room knows it. The one thing worth "
                 "dwelling on is resumability: real borrowers abandon mid-KYC, and the session-at-step-one "
                 "design is why they come back to the right screen instead of starting over. If you're "
                 "driving the live app, this is where you show the OTP arriving on a handset.",
    },
    {
        "kind": "shot",
        "kicker": "borrower journey  ·  2 of 5",
        "title": "Identity: Aadhaar via DigiLocker, then a liveness selfie",
        "sub": "This is the step most demos fake. Ours is a live provider journey, end to end.",
        "shot": "borrower-identity",
        "frame": "DigiLocker consent hand-off  +  liveness video capture",
        "route": "/signup/digilocker  »  /kyc/digilocker/callback  ·  /signup/selfie",
        "body": [
            "DigiLocker: the borrower authorises in DigiLocker's own flow. We accept the Aadhaar only "
            "when the document's digital signature validates.",
            "We pull the Aadhaar photo and store it — that photo becomes the reference image for the "
            "next step.",
            "Selfie: an interactive passive-liveness video, face-matched 1:1 against the Aadhaar photo. "
            "Not an uploaded photo of a photo.",
            "If the provider's liveness journey is unavailable, it degrades to a synchronous face-match "
            "fallback rather than blocking the borrower.",
            "> Neither path ever hard-rejects. A human KYC approver makes the final call — with the "
            "machine's verdict in front of them.",
        ],
        "notes": "This is the slide that separates us from a demo build, so slow down. Two specifics "
                 "to say out loud: (1) we gate on the Aadhaar document's digital signature, so a "
                 "screenshot of an Aadhaar is worthless; (2) the selfie is matched against the "
                 "DigiLocker photo, not against a self-uploaded ID. Then the operational point: it "
                 "degrades instead of blocking, because a provider outage must not stop origination.",
    },
    {
        "kind": "shot",
        "kicker": "borrower journey  ·  3 of 5",
        "title": "Underwriting inputs: bureau, salary, penny-drop",
        "sub": "Where the limit comes from — and where the money is allowed to land.",
        "shot": "borrower-underwriting",
        "frame": "Bureau pull  ·  salary verification  ·  penny-drop bank check",
        "route": "/signup/bureau  ·  /signup/salary  ·  /signup/penny-drop",
        "body": [
            "Bureau: a full report, not just a score. Balances, accounts and enquiries are parsed into "
            "structured facts.",
            "Those facts produce a 1–5star recommendation and a one-page branded credit brief PDF for "
            "staff — never shown to the borrower.",
            "Salary drives the eligible limit: 25% of verified monthly salary. Change the salary later "
            "and the limit recomputes, audited.",
            "Penny-drop confirms the bank account belongs to this person before anything is disbursed "
            "into it.",
            "> Two bureau providers are wired with automatic fallback, because a thin file at one "
            "bureau should not end the application.",
        ],
        "notes": "The commercial point here is fallback: a single-bureau integration means every thin "
                 "file is a lost customer. We route per capability with a second provider behind it. "
                 "Second point: penny-drop before disbursal is the cheapest fraud control in the "
                 "stack — one rupee to confirm the account name. Don't read the star rating as a "
                 "credit decision; it's a staff aid, and slide 16 says who actually decides.",
    },
    {
        "kind": "shot",
        "kicker": "borrower journey  ·  4 of 5",
        "title": "The limit, the amount, the true cost — before committing",
        "sub": "The borrower states their salary date, picks an amount, and sees the exact due date and total.",
        "shot": "borrower-apply",
        "frame": "Amount chooser with live due date + cost breakdown",
        "route": "/loan/apply",
        "body": [
            "Salary day is a single day-of-month field. The helper text previews the salary-linked due "
            "date using the same rule the backend will apply.",
            "Move the amount and fee, GST, net disbursal and total repayable update live. No surprise "
            "deductions on disbursal day.",
            "The amount is bounded server-side too: at least ₹1,000, never above the eligible limit.",
            "Repeat borrowers never re-state their salary day — it carries over from the prior loan.",
            "> The frontend math mirrors the backend engine exactly, so the quote the borrower agreed "
            "to is the quote the ledger books.",
        ],
        "notes": "Two things sell this screen. First, transparency: fee and GST are visible before "
                 "consent, which is where most complaints and most regulatory attention land. Second, "
                 "the salary-linked due date — this is the product's actual risk control, and the "
                 "borrower can see it move as they change the amount. If asked 'why not EMIs': one "
                 "instalment on payday means nothing to service and nothing to reconcile.",
    },
    {
        "kind": "shot",
        "kicker": "borrower journey  ·  5 of 5",
        "title": "Track, repay, borrow again",
        "sub": "The borrower sees their own audit trail, pays a prepayment-aware amount, and returns on a fast path.",
        "shot": "borrower-repay",
        "frame": "Live status trail  ·  repay screen  ·  borrow-again pre-approval",
        "route": "/loan/status  ·  /dashboard  ·  /repay  ·  /reloan",
        "body": [
            "Status is the real state machine, with the audit trail rendered — not a progress bar that "
            "guesses.",
            "Repay shows 'pay today', interest charged only to the day of payment. Prepayment is a "
            "first-class case, not an exception.",
            "The borrower records the payment; an accountant verifies it. The loan closes only when the "
            "penalty-aware balance actually reaches zero.",
            "Borrow again: a clean repayment history skips KYC and credit entirely and goes straight to "
            "disbursement. Any past overdue routes to a human review queue first.",
            "> One live loan at a time, with top-up allowed against remaining headroom.",
        ],
        "notes": "The reborrow path is the revenue slide of the borrower half — repeat borrowers are "
                 "the margin, and here they get money in minutes without re-verification, while "
                 "anyone who was ever overdue is quietly routed to a human. Say 'past delinquency is "
                 "the only gate — not the credit score', because that's a deliberate product choice "
                 "the client made and it usually starts a good conversation.",
    },
    {
        "kind": "section",
        "num": "02",
        "title": "The back office",
        "sub": "Nine roles, one state machine, and an audit trail that survives a regulator. "
               "This is the half that vendors leave you to build.",
        "notes": "Transition beat. Say: 'Everything so far was the easy half. Any agency can build a "
                 "pretty onboarding funnel. What follows is the part that decides whether you can "
                 "actually operate — and pass an audit.'",
    },
    {
        "kind": "chain",
        "kicker": "the spine",
        "title": "One status field walks the entire lifecycle",
        "sub": "Every transition is validated against a transition map, checked against the actor's role, and appended to an immutable event log.",
        "per_row": 5,
        "steps": [
            ("KYC_PENDING", "Borrower submits"), ("KYC_APPROVED", "KYC approver"),
            ("CREDIT_EXEC", "Executive recommends"), ("CREDIT_HEAD", "Head approves — SoD checked"),
            ("DISBURSEMENT", "Head releases"), ("ACCOUNTANT", "Transfer validated"),
            ("*ACTIVE", "Loan minted, due date set"), ("OVERDUE", "Penalty accrues, capped"),
            ("COLLECTIONS", "Case, DPD, settlement"), ("CLOSED", "Balance reaches zero"),
        ],
        "foot": "No stage-skipping: an invalid jump is refused with ILLEGAL_TRANSITION, a wrong role with "
                "FORBIDDEN_ROLE, a same-actor approval with SOD_VIOLATION. Rejection and cancellation paths "
                "exist at every stage. The event log is append-only and is itself the source of truth for the SoD check.",
        "notes": "This is the architectural claim of the whole deck: one row, one status, one log. "
                 "Contrast it explicitly — 'in a three-vendor stack, this diagram is spread across "
                 "three databases and nobody owns the join.' Name the three error codes; they're "
                 "concrete proof the rules are server-side, not UI convention.",
    },
    {
        "kind": "table",
        "kicker": "the back office",
        "title": "Nine roles, and exactly what each one may do",
        "sub": "Permissions are enforced in the services off a signed token — never in the browser, never in middleware.",
        "head": ["Role", "Owns this decision"],
        "widths": [0.28, 0.72],
        "rows": [
            ["KYC Approver", "Approves or rejects KYC; clears returning borrowers who were ever overdue; can wave a low-risk file straight to disbursement."],
            ["Credit Executive", "Recommends or rejects — the maker half of the credit decision."],
            ["Credit Head", "Assigns the executive and gives final approval — the checker half. Cannot approve their own recommendation."],
            ["Disbursement Head", "Releases funds: straight through with a transaction id, or routed to the accountant. Retries failures. Settles referral payouts."],
            ["Accountant", "Validates the outgoing transfer (this is what mints the loan) and verifies or rejects borrower repayments."],
            ["Collection Head", "Runs collections; approves or rejects settlements proposed by their own team."],
            ["Collection Executive", "Works cases and logs interactions. Salary and employer fields are hidden from this role."],
            ["Admin", "Oversight across every queue; audited KYC and salary corrections; blocklist, expenses, exports."],
            ["Developer", "Read-only. Can see customers, can decide nothing."],
        ],
        "foot": "Borrower and staff sessions are separate namespaces — different login endpoints, different cookies, "
                "different token audiences. A borrower token cannot address a staff endpoint even in principle.",
        "notes": "Read two rows only: Credit Head ('cannot approve their own recommendation') and "
                 "Collection Executive ('salary and employer hidden'). Those two show the granularity "
                 "without reading a table aloud. The footer is for the security-minded buyer: two "
                 "separate token audiences, not one system with a role flag.",
    },
    {
        "kind": "shot",
        "kicker": "back office  ·  1 of 8",
        "title": "The console: one workbench, your queues only",
        "sub": "Every staff member signs into the same page and sees the work that is actually theirs.",
        "shot": "staff-dashboard",
        "frame": "Staff dashboard + live application pipeline",
        "route": "/staff/dashboard  ·  /staff/applications",
        "body": [
            "The dashboard opens on 'N items need your action' for that role — a personal queue, not a "
            "report to interpret.",
            "Thirty-day sparklines for applications, disbursals and repayments, plus a pipeline bar "
            "across every stage.",
            "The pipeline page is a literal rendering of the state machine: one panel per status, each "
            "with its stage's action buttons.",
            "Admin sees all eight panels at once as an oversight surface; every other role sees only "
            "the panels it can act on.",
            "> The closed-loans archive is deliberately lazy-loaded, so the console never polls "
            "history nobody is looking at.",
        ],
        "notes": "The operational sell: no training manual. A new hire signs in and the page already "
                 "contains only what they may do. When you demo this live, switch roles once and let "
                 "them watch the panels change — that single moment does more than this slide. Let "
                 "the page settle before you narrate; there's a brief placeholder while the session "
                 "resolves.",
    },
    {
        "kind": "shot",
        "kicker": "back office  ·  2 of 8",
        "title": "Verification dashboard: triage, not a spreadsheet",
        "sub": "Which borrowers are stuck, on what, and whose problem it is.",
        "shot": "staff-verifications",
        "frame": "Verification dashboard — tiles, triage buckets, manual override dialog",
        "route": "/staff/verifications",
        "body": [
            "Five tiles: pending, failed, in review, passed, never run — across every in-flight "
            "application.",
            "Four triage buckets in priority order: has failures, awaiting the borrower, all passed, "
            "not started.",
            "Open any applicant for the check-by-check view with each provider's verdict.",
            "Manual override: a KYC approver can pass or fail a single check with remarks — recorded as "
            "a MANUAL decision, in the audit trail.",
            "> One-click reminder nudges the borrower on exactly the steps still outstanding.",
        ],
        "notes": "Say the honest thing about overrides: automated KYC is never 100%, so the choice is "
                 "between an audited override and staff finding a workaround. We built the audited "
                 "override. Note the dashboard tracks pending work only — decided applications leave "
                 "it, by design; it's a work queue, not a history log. Don't hunt for an approved "
                 "applicant here on camera.",
    },
    {
        "kind": "shot",
        "kicker": "back office  ·  3 of 8",
        "title": "KYC clearance — and the deliberate fast path",
        "sub": "Two queues that are kept permanently apart: new files, and returning borrowers with history.",
        "shot": "staff-kyc",
        "frame": "KYC approvals queue  +  reborrow review queue",
        "route": "/staff/kyc-approvals  ·  /staff/kyc-review",
        "body": [
            "New applications: approve or reject, with the full verification picture on the same screen.",
            "Below it, the instant-loan panel — a KYC approver can send a low-risk file straight to the "
            "Disbursement Head, skipping the credit chain.",
            "Separate queue: returning borrowers who were ever overdue, each row showing their loan "
            "history inline.",
            "Buttons there read 'Clear borrower' or 'Reject' — different decision, different language.",
            "> A clean repayment history never reaches this queue at all. Only past delinquency does.",
        ],
        "notes": "The fast path is a policy choice made explicit in software: your credit team's time "
                 "goes to the files that need it. Be clear that it's configurable policy, not a "
                 "shortcut we took — the full maker-checker chain is one setting away, and the "
                 "unreviewed fast path is still fully audited. Prospects with a conservative credit "
                 "committee will ask exactly this.",
    },
    {
        "kind": "shot",
        "kicker": "back office  ·  4 of 8",
        "title": "Credit: assign, recommend, approve",
        "sub": "The maker-checker chain, with the bureau brief on the same screen as the decision.",
        "shot": "staff-credit",
        "frame": "Application detail — credit badge, journey stepper, cost card, audit log",
        "route": "/staff/applications » Open  ·  /staff/credit/queue",
        "body": [
            "The head assigns an executive — and the picker only lists currently active executives, so "
            "work is never assigned into a void.",
            "Application detail carries a staff-only score and star rating pill, plus the one-page bureau "
            "brief as a downloadable PDF.",
            "The journey stepper shows every step with its evidence; the cost card shows fee, GST, net "
            "disbursal and total repayable.",
            "A collapsed audit log lists every action, actor and timestamp on that application.",
            "> The score is never shown to the borrower, and risk band never changes the price. It "
            "informs a human decision; it does not make one.",
        ],
        "notes": "Point at the score pill and say the two rules together: staff-only, and it does not "
                 "move the price. Then open the audit log — 'every action, actor, timestamp' is the "
                 "line that lands with anyone who has sat through an audit. If you land on a "
                 "fast-tracked reborrow, call out the fast-track marker so nobody thinks a stage was "
                 "skipped silently.",
    },
    {
        "kind": "bullets",
        "kicker": "the differentiator",
        "title": "Separation of duties, enforced by replaying the audit log",
        "sub": "Not a UI convention. Not a config flag. The server refuses.",
        "cols": [[
            "## How it actually works",
            "Every maker-checker approval replays that application's event log before it commits.",
            "If the actor now approving is the same staff id that recommended, the call is rejected with "
            "SOD_VIOLATION. There is no UI path around it, because the check is not in the UI.",
            "The same shape applies to settlements — the officer who proposed one cannot approve it — "
            "and to repayment verification.",
            "## Every approve has a paired reject",
            "Credit, disbursement, settlements, repayment verification: each approval has a real, "
            "audited counter-action. We refuse to ship approve-only flows, because that is how staff "
            "end up faking approvals to clear a queue.",
        ], [
            "## Why this is the slide that matters",
            "This is the control an auditor asks about first, and the one bolted-together stacks cannot "
            "produce — because no single system saw both halves of the decision.",
            "Because the log is append-only and is itself the input to the check, the control cannot "
            "drift out of sync with the record of it.",
            "## The one documented exception",
            "Admin is deliberately exempt, so oversight can walk a stuck loan through one step at a "
            "time. It is a written product decision, and every step it takes is still logged under the "
            "admin's own staff id.",
            "> We will show you the exemption on screen too. A control you cannot describe the limits "
            "of is not a control.",
        ]],
        "notes": "Slow down and stop moving the mouse. This is the whole pitch in one slide. The "
                 "strongest move in the room is volunteering the admin exemption before anyone asks — "
                 "it converts 'vendor claims compliance' into 'these people know where their own "
                 "edges are.' If you demo it, use two real personas handing off; do not try to "
                 "trigger the violation live as admin, because admin is exactly the exempt case.",
    },
    {
        "kind": "shot",
        "kicker": "back office  ·  5 of 8",
        "title": "Disbursement and accounting: money moves on a human confirmation",
        "sub": "Two release paths, one ledger, no automatic reconciliation.",
        "shot": "staff-disbursement",
        "frame": "Disbursement panels  ·  transfers to confirm  ·  transactions ledger",
        "route": "/staff/applications (Disb. Head, Accountant)  ·  /staff/accounting/transactions",
        "body": [
            "Fast release: the Disbursement Head enters the transaction id and the loan activates "
            "immediately, reference recorded.",
            "Controlled release: approve without an id and it routes to the Accountant to validate the "
            "transfer first.",
            "Confirming that transfer is what mints the loan and sets the salary-linked due date. "
            "Failures land in a retry panel, not a void.",
            "Repayments the borrower reports queue for accountant verification — verifying is what "
            "reduces the balance and can close the loan.",
            "> One company-wide ledger of every disbursal out and repayment in, searchable, date-ranged, "
            "exportable.",
        ],
        "notes": "The design principle worth stating: no automatic reconciliation. A human confirms "
                 "that money actually moved, because a false positive here is a real loss. Two "
                 "release paths exist because a two-person control on every disbursal is right at "
                 "scale and paralysing at ten loans a day — the client chose per-case. On the ledger, "
                 "set the period to 'all time' before you talk, or it will look empty.",
    },
    {
        "kind": "shot",
        "kicker": "back office  ·  6 of 8",
        "title": "Collections: buckets, cases, and arithmetic you can read",
        "sub": "Assign from the row, log every contact with proof, and settle under maker-checker.",
        "shot": "staff-collections",
        "frame": "Overdue queue with live DPD  ·  DPD buckets  ·  case detail  ·  settlements",
        "route": "/staff/applications  ·  /staff/collections/settlements",
        "body": [
            "The overdue queue shows live days-past-due per loan; the head assigns an officer straight "
            "from the row, which opens the case.",
            "DPD buckets — upcoming, 1–7, 8–30, 31–60, 61–90, 90+ — computed on read, never a stale "
            "stored field.",
            "The amount due card spells out its own arithmetic: 1%/day interest across N days, 2%/day "
            "penalty capped at 30. Auditable, not a black box.",
            "Interaction log demands a proof reference the moment an outcome is recorded as PAID — "
            "otherwise the API refuses.",
            "> Settlements are propose-then-approve: the officer who proposed one can never approve it.",
        ],
        "notes": "Two lines to say: 'computed on read, never stored' (so DPD cannot silently rot) and "
                 "'proof is mandatory on PAID' (so a collections officer cannot clear a case on their "
                 "word alone). Settlement maker-checker is the same control as credit, applied where "
                 "the write-off risk actually is. Note the executive cannot see salary or employer "
                 "fields — need-to-know, enforced server-side.",
    },
    {
        "kind": "shot",
        "kicker": "back office  ·  7 of 8",
        "title": "Customers and administration",
        "sub": "One borrower-centric view, and the levers an operations lead actually asks for on day two.",
        "shot": "staff-admin",
        "frame": "Customer 360  ·  staff roster  ·  blocklist  ·  payment settings  ·  expenses",
        "route": "/staff/customers  ·  /staff/admin/*",
        "body": [
            "Customer 360: every application, loan, payment and profile change in one place, searchable "
            "by name or id.",
            "Admin corrections are audited field by field, previous value beside new. Identity fields "
            "(PAN, Aadhaar, mobile) are locked.",
            "A salary correction recomputes the eligible limit automatically — the derived number can "
            "never contradict the stored one.",
            "Fraud blocklist by PAN, Aadhaar reference, phone, device or bank account. Staff roster, "
            "one-time invites, payee details, expense ledger, full application register.",
            "> Data export is admin-only across the entire console — a deliberate governance decision, "
            "not a per-page setting.",
        ],
        "notes": "The audited-edit beat demos beautifully: bump a salary, then scroll to change history "
                 "and show the row that just appeared. Do it live if you can. The admin-only export "
                 "rule is worth naming — it's the answer to 'how do I stop a junior walking out with "
                 "the book?' Everything on this slide is the stuff prospects discover they need in "
                 "month two and then pay to retrofit.",
    },
    {
        "kind": "bullets",
        "kicker": "how it is built",
        "title": "Architecture: boring on purpose",
        "sub": "Java 21 / Spring Boot 3 and Next.js 15, deployed on managed AWS. Nothing exotic to hire for.",
        "cols": [[
            "## The shape",
            "Next.js front end with a backend-for-frontend layer — the browser never talks to the core "
            "API directly, and no secret ever reaches it.",
            "Spring Boot core, split into modules by domain: loan, KYC, verification, collections, "
            "notifications, IAM, storage.",
            "One bootable service. Postgres for state, S3 for documents, SSM for secrets.",
            "Deployed Vercel » ALB » ECS Fargate » RDS, with a written redeploy runbook and smoke tests.",
        ], [
            "## The decisions that will still hold in year three",
            "Authorisation lives in the services, off a signed token — not in middleware, where it "
            "silently stops covering new routes.",
            "Business rules live in one place each: one money engine, one flow service, one transition map.",
            "Schema changes are versioned migrations applied on boot. No hand-run DDL, no drift between "
            "environments.",
            "Verification providers sit behind one capability-routed port, so swapping or adding one "
            "touches no business logic.",
            "> A team that knows Spring and React can own this. That is a hiring decision as much as a "
            "technical one.",
        ]],
        "notes": "Aim this at the CTO or the technical advisor in the room. The closing line is the "
                 "real point: we chose a stack they can hire for, not one that keeps them dependent "
                 "on us. Saying that out loud tends to win the technical veto. If they push for "
                 "detail, the deployment runbook and every resource id are documented in the repo — "
                 "offer to walk it after.",
    },
    {
        "kind": "table",
        "kicker": "how it is built",
        "title": "Integrations already shipped and live-tested",
        "sub": "Each one sits behind an internal port, so a provider is a configuration change rather than a rewrite.",
        "head": ["Capability", "How it is wired", "State"],
        "widths": [0.26, 0.5, 0.24],
        "rows": [
            ["PAN, bureau, email", "Primary provider with an automatic second-provider fallback per capability", "Live-tested"],
            ["Aadhaar (DigiLocker)", "Consent hand-off, signature-validated, Aadhaar photo ingested to storage", "Live-tested"],
            ["Selfie liveness", "Interactive video journey, 1:1 match to the Aadhaar photo, face-match fallback", "Live-tested"],
            ["Penny-drop", "Bank account ownership confirmed before any disbursal", "Live-tested"],
            ["Address", "Provider address-verification API", "Live-tested"],
            ["Email delivery", "AWS SES, branded HTML, bounce and complaint feedback auto-suppressed", "Live, sandbox"],
            ["SMS", "DLT-compliant gateway; 15 templates registered, entity and sender approved", "Templates in review"],
            ["Documents", "S3 with presigned access; generated agreements and credit briefs stored", "Live"],
        ],
        "foot": "The bounce loop is worth a sentence: a hard bounce or a complaint adds that address to a suppression "
                "list automatically, and the sender skips it forever after. Sender reputation is protected without anyone watching a dashboard.",
        "notes": "Do not read this table. Say: 'every one of these is live-tested against the real "
                 "provider, and each sits behind a port so switching vendors is configuration.' Then "
                 "tell the bounce-suppression story — it's small, concrete, and signals the kind of "
                 "thing that only shows up in software written by people who have operated it. Be "
                 "straight about the two 'in review' rows; slide 25 covers them.",
    },
    {
        "kind": "bullets",
        "kicker": "how it is built",
        "title": "The engineering decisions that keep the numbers honest",
        "sub": "Money and audit are where lending systems quietly rot. These four choices are why this one does not.",
        "cols": [[
            "## Integer paise, everywhere",
            "Every amount is a whole number of paise with explicit rounding. No floating point touches "
            "money at any layer, front end included.",
            "## One balance, computed on read",
            "Outstanding is derived every time from principal, interest to date, capped penalty and "
            "verified payments — the stored column is only a cache.",
            "Every surface that shows 'amount owed' calls the same function, so the borrower's screen, "
            "the collections queue and the ledger can never disagree.",
            "A loan closes only when that penalty-aware balance truly reaches zero.",
        ], [
            "## Append-only audit as a load-bearing structure",
            "The event log is not a report; it is the input to the separation-of-duties check. It "
            "cannot be stale, because the control reads it.",
            "## Notifications that cannot break a loan",
            "Domain code publishes an event and commits. Delivery happens after, asynchronously. A "
            "dead SMS gateway can never fail a disbursal.",
            "In-app, SMS and email from one engine, with per-borrower preferences honoured.",
            "> Plus database-backed feature flags, so a feature can be switched off in production "
            "without a redeploy.",
        ]],
        "notes": "For a finance audience, 'no floating point touches money' and 'one balance function, "
                 "everywhere' are the two sentences that matter — most of them have been burned by a "
                 "system where two screens disagreed about what a customer owed. For a technical "
                 "audience, the audit log being the input to the control (not a byproduct) is the "
                 "elegant bit. Pick your half based on who is in the room.",
    },
    {
        "kind": "stats",
        "kicker": "how it is handed over",
        "title": "What ships alongside the software",
        "sub": "A platform you cannot operate or verify is not delivered. This is what makes it yours.",
        "stats": [
            ("51 classes", "380+ backend test methods covering the money engine, the flow rules and SoD."),
            ("Integration suite", "The full lifecycle exercised against a real Postgres, in a container."),
            ("7 e2e specs", "Playwright over auth, RBAC, onboarding, repay, staff console and admin."),
            ("CI on every push", "Build, tests and type-checks run automatically. Nothing merges unverified."),
            ("Seeded demo env", "One command: 49 applications at every stage, plus collections and back office."),
            ("Onboarding doc", "A single document that takes a new engineer from clone to running stack."),
            ("Deploy runbook", "Every cloud resource, the redeploy recipe and smoke tests, written down."),
            ("Walkthrough script", "A chaptered demo script — your sales team can run this pitch without us."),
        ],
        "foot": "The last one is not a throwaway: we wrote the script that a non-technical presenter uses to demo "
                "the product, including which screens look empty and why. Documentation is part of the deliverable, not a follow-up invoice.",
        "notes": "This slide answers the fear nobody says out loud: 'will we be hostage to this "
                 "vendor?' Lead with the seeded demo environment and the onboarding doc — those are "
                 "what let their own team take over. The walkthrough script line is unusual enough "
                 "that it tends to get remembered.",
    },
    {
        "kind": "bullets",
        "kicker": "straight talk",
        "title": "What is not finished — and what it takes",
        "sub": "We would rather you hear this from us in slide 25 than find it in week three.",
        "cols": [[
            "## In flight, waiting on third parties",
            "SMS templates: entity and sender are approved; a re-filed template batch is in operator "
            "review. Until they clear, transactional SMS is off. Not code — telecom queue time.",
            "Email: fully working, still inside the provider's sandbox. Production access is an "
            "application, and the suppression handling it needs is already built.",
            "## Needs a build, scoped and understood",
            "Real bank payout at the accountant step: today the transfer is made in the bank and the "
            "reference recorded. Direct NEFT/IMPS integration is a defined next phase.",
        ], [
            "## Hardening before high volume",
            "Foreign-key constraints and encryption of personal data at rest — both known, both "
            "sequenced, neither blocking a controlled launch.",
            "The daily reminder job runs on a single instance. It needs a distributed lock before the "
            "service scales horizontally.",
            "## How we would sequence it with you",
            "One phase, a few weeks, run in parallel with your licensing and compliance sign-off — "
            "because those calendars are longer than this list anyway.",
            "> Everything above is written down in the repository's own readiness checklist. We did "
            "not assemble this slide for you; we kept it as we built.",
        ]],
        "notes": "Counter-intuitively this is the slide that closes deals. Deliver it flatly, no "
                 "hedging, no apology. The last line is the strongest thing on it: this list already "
                 "existed as an internal checklist, which means we track our own gaps as a habit. If "
                 "someone asks 'so it's not production-ready?' — the honest answer is that the "
                 "software is; the third-party approvals and the payout integration are the runway, "
                 "and both are shorter than an NBFC's own licensing timeline.",
    },
    {
        "kind": "table",
        "kicker": "working with us",
        "title": "Three ways to start",
        "sub": "Same engineering team, same spine, different amount of your product in it.",
        "head": ["Engagement", "What you get", "Indicative timeline", "Commercials"],
        "widths": [0.2, 0.42, 0.19, 0.19],
        "rows": [
            ["White-label", "This platform, your brand, your rules configured: limits, fee, tenure, roles, workflow. Your providers wired in.", "Weeks, not months", "TODO — insert pricing"],
            ["Build on the spine", "Your product — different collateral, tenure, EMI structure or segment — on the same state machine, audit and SoD core.", "One quarter, typically", "TODO — insert pricing"],
            ["Team augmentation / AMC", "Our engineers inside your team, or we run and maintain what is already live. Documentation and handover included.", "Ongoing, monthly", "TODO — insert pricing"],
        ],
        "foot": "TODO before presenting: fill in commercials, and decide whether to name DhanBoost as the reference "
                "client on this slide or keep it anonymous until the room is qualified.",
        "notes": "Do not guess pricing on your feet — the placeholders are there so you fill them in "
                 "before the meeting, per prospect. Ask which of the three they think they are and "
                 "let them self-select; the answer tells you the deal size faster than any qualifying "
                 "question. The white-label timeline is the hook: 'weeks' is only credible because "
                 "they just watched the product work.",
    },
    {
        "kind": "close",
        "title": "The difference between a demo\nand a platform is the audit trail.",
        "body": [
            "Everything in this deck is running software. The screens are real, the providers are real, "
            "the state machine refuses invalid moves in production today.",
            "We build the half that is hard to see: separation of duties, one honest balance, an "
            "append-only record, and documentation that lets your team take over.",
            "Next step — a live driven demo on your questions, not our script. Bring your credit "
            "policy and your compliance lead; we will walk a loan from sign-up to closure and stop "
            "wherever you want to look underneath.",
        ],
        "foot": "SoftSolutionsAI  ·  softsolutionsai.com  ·  TODO: contact name, email, phone",
        "notes": "Close on the invitation, not a summary. 'Bring your compliance lead' is the "
                 "confident move — vendors who fake it never say that. Fill in the contact line "
                 "before you present, and if the deck is being left behind, add the reference-client "
                 "line you agreed on slide 27.",
    },
]


# ---------------------------------------------------------------- build


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for d in SLIDES:
        slide = BUILDERS[d["kind"]](prs, d)
        slide.notes_slide.notes_text_frame.text = d["notes"]
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)

    # self-check: reopen what we wrote and assert the deck is actually presentable
    check = Presentation(OUT)
    assert len(check.slides) == len(SLIDES), f"{len(check.slides)} slides != {len(SLIDES)} defined"
    for i, slide in enumerate(check.slides, 1):
        assert slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip(), \
            f"slide {i} has no speaker notes"
    spills = []
    for tb, paras, w, h in OVERFLOW:
        used = sum(
            est_lines(body, w / Inches(1), size) * size * 1.15 / 72 + after / 72
            for body, size, _c, _b, after in paras
        )
        avail = h / Inches(1)
        if used > avail * 1.02:
            spills.append(f"    {paras[0][0][:52]!r} needs ~{used:.2f}in, has {avail:.2f}in")
    if spills:
        print(f"WARNING — {len(spills)} text block(s) may overflow their box:")
        print("\n".join(spills))

    missing = [d["shot"] for d in SLIDES if d["kind"] == "shot" and not d.get("_shot_ok")]
    print(f"{OUT.relative_to(REPO)}  ·  {len(check.slides)} slides, notes on all")
    if missing:
        print(f"placeholder frames awaiting screenshots ({len(missing)}): {', '.join(missing)}")
        print("  » capture per docs/pitch/SHOTLIST.md into docs/pitch/shots/, then re-run this script")
    else:
        print("all screenshots embedded")


if __name__ == "__main__":
    main()
